#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=""):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

    return slide

def create_section_header_slide(prs, title):
    """섹션 헤더 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_shape = slide.shapes.title
    title_shape.text = title

    # 배경색 변경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel Blue

    return slide

def create_content_slide(prs, title, content_items):
    """컨텐츠 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title

    # 본문 영역
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('font_size', 18))

    return slide

def create_class_detail_slide(prs, class_name, package, role, why, benefits, alternatives):
    """클래스 상세 설명 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 타이틀
    title_shape = slide.shapes.title
    title_shape.text = f"{class_name}"

    # 본문
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # 패키지
    p = tf.paragraphs[0]
    p.text = f"📦 패키지: {package}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(10)

    # 역할
    p = tf.add_paragraph()
    p.text = f"🎯 역할"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = role
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)

    # 사용 이유
    p = tf.add_paragraph()
    p.text = f"💡 사용 이유"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 76)

    p = tf.add_paragraph()
    p.text = why
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)

    # 장점
    p = tf.add_paragraph()
    p.text = f"✅ 주요 장점"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 102, 0)

    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = f"• {benefit}"
        p.font.size = Pt(13)
        p.level = 1

    if alternatives:
        p.space_after = Pt(10)

        # 대안
        p = tf.add_paragraph()
        p.text = f"🔄 대안"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(153, 0, 153)

        p = tf.add_paragraph()
        p.text = alternatives
        p.font.size = Pt(13)
        p.level = 1

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    create_title_slide(prs,
                      "Hoji Backend Base Project",
                      "Spring Boot + Kotlin 기반 백엔드 베이스 프로젝트 상세 가이드")

    # 2. 목차
    create_content_slide(prs, "목차", [
        {'text': '1. 프로젝트 개요', 'level': 0, 'font_size': 20},
        {'text': '2. 아키텍처 구조', 'level': 0, 'font_size': 20},
        {'text': '3. 공통 모듈 (Common)', 'level': 0, 'font_size': 20},
        {'text': 'DTO, Exception, Logging, Context, Metrics, Health, Util', 'level': 1, 'font_size': 16},
        {'text': '4. 설정 모듈 (Config)', 'level': 0, 'font_size': 20},
        {'text': 'JPA, QueryDSL, Redis, RabbitMQ, WebClient, Async, CORS 등', 'level': 1, 'font_size': 16},
        {'text': '5. 도메인 모듈 (Domain)', 'level': 0, 'font_size': 20},
        {'text': '6. 레포지토리, 서비스, 컨트롤러', 'level': 0, 'font_size': 20},
        {'text': '7. 운영 고려사항', 'level': 0, 'font_size': 20},
    ])

    # 3. 프로젝트 개요
    create_section_header_slide(prs, "프로젝트 개요")

    create_content_slide(prs, "기술 스택", [
        {'text': '코어 기술', 'level': 0, 'font_size': 20},
        {'text': 'Java 17 + Kotlin 1.9.20', 'level': 1, 'font_size': 16},
        {'text': 'Spring Boot 3.2.0', 'level': 1, 'font_size': 16},
        {'text': 'Gradle 8.5 (Kotlin DSL)', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 14},
        {'text': 'Database & ORM', 'level': 0, 'font_size': 20},
        {'text': 'H2 (In-Memory)', 'level': 1, 'font_size': 16},
        {'text': 'Spring Data JPA + QueryDSL', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 14},
        {'text': 'Infra & Monitoring', 'level': 0, 'font_size': 20},
        {'text': 'Redis, RabbitMQ', 'level': 1, 'font_size': 16},
        {'text': 'Micrometer + Prometheus', 'level': 1, 'font_size': 16},
        {'text': 'Logback', 'level': 1, 'font_size': 16},
    ])

    create_content_slide(prs, "프로젝트 목표", [
        {'text': '상용 서비스를 위한 견고한 기반 제공', 'level': 0, 'font_size': 18},
        {'text': '일관된 응답 구조 및 예외 처리', 'level': 1, 'font_size': 16},
        {'text': '포괄적인 로깅 및 모니터링', 'level': 1, 'font_size': 16},
        {'text': '멀티 인스턴스 환경 지원', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 14},
        {'text': '확장 가능하고 유지보수가 용이한 구조', 'level': 0, 'font_size': 18},
        {'text': '모듈화된 패키지 구조', 'level': 1, 'font_size': 16},
        {'text': '환경별 설정 분리 (Local, Dev, Prod)', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 14},
        {'text': '운영 편의성', 'level': 0, 'font_size': 18},
        {'text': 'Health Check, Graceful Shutdown', 'level': 1, 'font_size': 16},
        {'text': 'JPA Auditing, Async 처리', 'level': 1, 'font_size': 16},
    ])

    # 4. 아키텍처
    create_section_header_slide(prs, "아키텍처 구조")

    create_content_slide(prs, "패키지 구조", [
        {'text': 'common/ - 공통 기능 모듈', 'level': 0, 'font_size': 18},
        {'text': 'dto, exception, logging, context, metrics, health, util', 'level': 1, 'font_size': 14},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'config/ - 설정 모듈', 'level': 0, 'font_size': 18},
        {'text': 'JPA, QueryDSL, Redis, RabbitMQ, WebClient, Async, CORS, Metrics, Swagger', 'level': 1, 'font_size': 14},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'domain/ - 도메인 엔티티', 'level': 0, 'font_size': 18},
        {'text': 'BaseEntity (공통), User (예제)', 'level': 1, 'font_size': 14},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'repository/ - 데이터 액세스', 'level': 0, 'font_size': 18},
        {'text': 'service/ - 비즈니스 로직', 'level': 0, 'font_size': 18},
        {'text': 'controller/ - API 엔드포인트', 'level': 0, 'font_size': 18},
    ])

    # 5. Common - DTO
    create_section_header_slide(prs, "Common: DTO")

    create_class_detail_slide(prs,
        "ApiResponse",
        "com.hoji.common.dto",
        "모든 API 응답을 일관된 구조로 래핑하는 제네릭 클래스",
        "API 클라이언트가 모든 응답을 동일한 형태로 받아 처리할 수 있도록 표준화",
        [
            "일관된 응답 구조로 클라이언트 개발 용이",
            "성공/실패 여부, 에러 코드, 메시지를 표준화",
            "타임스탬프로 응답 시간 추적 가능",
            "제네릭을 통한 타입 안정성"
        ],
        "Spring HATEOAS, Problem Details (RFC 7807) 등이 있으나, 단순하고 직관적인 구조로 충분한 경우 현재 방식이 효율적"
    )

    create_class_detail_slide(prs,
        "ResultCode",
        "com.hoji.common.dto",
        "응답 코드와 메시지를 정의하는 Enum 클래스",
        "에러 코드를 중앙 집중식으로 관리하여 일관성 유지",
        [
            "코드와 메시지를 한 곳에서 관리",
            "오타나 중복 코드 방지",
            "다국어 지원 시 확장 용이",
            "Enum의 타입 안정성 활용"
        ],
        "외부 파일(YAML, properties) 관리도 가능하나, Enum 방식이 컴파일 타임 체크와 IDE 지원에 유리"
    )

    # 6. Common - Exception
    create_section_header_slide(prs, "Common: Exception")

    create_class_detail_slide(prs,
        "BusinessException",
        "com.hoji.common.exception",
        "비즈니스 로직에서 발생하는 예외를 표현하는 커스텀 예외 클래스",
        "시스템 예외와 비즈니스 예외를 명확히 구분하여 적절한 처리 가능",
        [
            "ResultCode와 연동하여 일관된 에러 응답",
            "스택 트레이스 제어로 성능 최적화 가능",
            "명시적인 예외 타입으로 가독성 향상",
            "체크 예외가 아닌 런타임 예외로 보일러플레이트 감소"
        ],
        "Vavr의 Try, Either 패턴도 고려 가능하나, Spring 환경에서는 예외 기반 처리가 더 자연스러움"
    )

    create_class_detail_slide(prs,
        "GlobalExceptionHandler",
        "com.hoji.common.exception",
        "@RestControllerAdvice를 사용한 전역 예외 처리기",
        "모든 컨트롤러에서 발생하는 예외를 한 곳에서 처리하여 일관된 에러 응답 제공",
        [
            "중복 예외 처리 코드 제거",
            "일관된 에러 응답 형식 보장",
            "환경별 에러 정보 노출 제어",
            "로깅과 예외 처리를 통합 관리"
        ],
        "Spring의 표준 기능이며, ControllerAdvice가 최선의 방법. @ExceptionHandler 조합으로 세밀한 제어 가능"
    )

    # 7. Common - Logging
    create_section_header_slide(prs, "Common: Logging")

    create_class_detail_slide(prs,
        "LoggingInterceptor",
        "com.hoji.common.logging",
        "모든 HTTP 요청/응답을 자동으로 로깅하는 인터셉터",
        "API 호출 추적, 디버깅, 성능 모니터링을 위한 상세 로그 제공",
        [
            "요청/응답 전체 정보 자동 로깅",
            "민감 정보(Authorization 등) 마스킹",
            "처리 시간 측정",
            "환경별 로깅 레벨 조정 가능"
        ],
        "Spring Cloud Sleuth + Zipkin (분산 추적), ELK Stack (중앙 로그 관리)과 조합 가능. 기본 로깅으로는 현재 방식이 적합"
    )

    create_class_detail_slide(prs,
        "RequestResponseCachingFilter",
        "com.hoji.common.logging",
        "요청/응답 본문을 여러 번 읽을 수 있도록 캐싱하는 필터",
        "Stream은 한 번만 읽을 수 있는데, 로깅과 실제 처리 모두에서 읽어야 하므로 캐싱 필요",
        [
            "요청/응답 본문을 중복 읽기 가능",
            "로깅 후에도 정상 처리 가능",
            "ContentCachingWrapper 활용"
        ],
        "필수적인 구현. 대안은 없으며 Spring이 제공하는 ContentCaching 클래스 활용이 표준"
    )

    # 8. Common - Context
    create_section_header_slide(prs, "Common: Context")

    create_class_detail_slide(prs,
        "RequestContext",
        "com.hoji.common.context",
        "ThreadLocal을 사용하여 요청별 컨텍스트 정보를 저장하는 유틸리티",
        "요청 ID, 사용자 ID 등을 서비스 레이어 깊숙이까지 파라미터 없이 전달",
        [
            "메서드 시그니처 오염 방지",
            "로깅, Auditing에서 사용자 정보 활용",
            "요청 추적을 위한 Request ID 관리",
            "메모리 누수 방지를 위한 cleanup 지원"
        ],
        "Spring Security의 SecurityContextHolder 패턴과 유사. 인증 정보는 SecurityContext, 기타 정보는 RequestContext 사용 권장"
    )

    create_class_detail_slide(prs,
        "RequestContextInterceptor",
        "com.hoji.common.context",
        "HTTP 헤더에서 컨텍스트 정보를 추출하여 RequestContext에 저장하는 인터셉터",
        "요청 시작 시 컨텍스트 초기화, 종료 시 cleanup 보장",
        [
            "자동으로 컨텍스트 생성/제거",
            "메모리 누수 방지",
            "X-Request-ID, X-User-ID 헤더 처리",
            "UUID 자동 생성으로 추적성 향상"
        ],
        "Servlet Filter로도 구현 가능하나, Interceptor가 Spring MVC 통합에 유리"
    )

    # 9. Common - Metrics
    create_section_header_slide(prs, "Common: Metrics")

    create_class_detail_slide(prs,
        "CustomMetrics",
        "com.hoji.common.metrics",
        "Micrometer를 래핑하여 비즈니스 메트릭을 쉽게 수집하는 유틸리티",
        "개발자가 메트릭 수집 코드를 간단하게 작성할 수 있도록 추상화",
        [
            "API 호출, 비즈니스 이벤트, 에러 카운팅",
            "처리 시간 측정 (time 메서드)",
            "Prometheus, Grafana와 연동",
            "멀티 인스턴스 환경에서 인스턴스별 추적"
        ],
        "Micrometer 직접 사용도 가능하나, 공통 패턴을 래핑하여 일관성과 생산성 향상"
    )

    create_class_detail_slide(prs,
        "MetricsInterceptor",
        "com.hoji.common.metrics",
        "모든 HTTP 요청에 대한 메트릭을 자동 수집하는 인터셉터",
        "개발자가 별도 코드 없이도 API 호출 통계를 자동으로 수집",
        [
            "API 엔드포인트별 호출 횟수",
            "HTTP 메서드, 상태 코드별 통계",
            "에러 발생 자동 추적",
            "처리 시간 측정"
        ],
        "Spring Boot Actuator의 기본 메트릭과 함께 사용. 커스텀 비즈니스 메트릭 추가"
    )

    # 10. Common - Health & Util
    create_section_header_slide(prs, "Common: Health & Util")

    create_class_detail_slide(prs,
        "CustomHealthIndicator",
        "com.hoji.common.health",
        "Spring Boot Actuator의 HealthIndicator 구현체",
        "Kubernetes liveness/readiness probe와 통합하여 무중단 배포 지원",
        [
            "표준 헬스 체크 인터페이스",
            "Kubernetes와 네이티브 통합",
            "여러 HealthIndicator 자동 조합",
            "모니터링 도구와 연동 용이"
        ],
        "Spring Boot Actuator 표준 방식. 커스텀 헬스 체크 로직 추가 가능"
    )

    create_class_detail_slide(prs,
        "DateTimeUtils & JsonUtils",
        "com.hoji.common.util",
        "날짜/시간 처리 및 JSON 변환을 위한 유틸리티 클래스",
        "반복적인 날짜 포맷팅, JSON 변환 코드를 중앙화하여 일관성 유지",
        [
            "날짜 포맷팅, 파싱, 계산 통합",
            "JSON 직렬화/역직렬화 간소화",
            "ObjectMapper 설정 통합 관리",
            "테스트 용이성"
        ],
        "Apache Commons, Joda Time 등도 있으나, Java 8+ LocalDateTime과 Jackson이면 충분"
    )

    # 11. Config 섹션
    create_section_header_slide(prs, "Config: 설정 모듈")

    create_class_detail_slide(prs,
        "JpaConfig",
        "com.hoji.config",
        "JPA 및 Hibernate 설정을 관리하는 설정 클래스",
        "EntityManager, TransactionManager 등 JPA 핵심 컴포넌트 설정",
        [
            "@EnableJpaRepositories로 레포지토리 스캔",
            "데이터소스 및 트랜잭션 관리",
            "Hibernate 프로퍼티 설정",
            "환경별 DDL 정책 제어"
        ],
        "Spring Boot의 자동 설정만으로도 충분하나, 세밀한 제어가 필요한 경우 명시적 설정 클래스 사용"
    )

    create_class_detail_slide(prs,
        "AuditConfig",
        "com.hoji.config",
        "JPA Auditing 설정 (@EnableJpaAuditing)",
        "엔티티의 생성/수정 정보를 자동으로 추적",
        [
            "createdBy, updatedBy 자동 설정",
            "RequestContext와 통합하여 사용자 정보 주입",
            "코드 중복 제거",
            "일관된 감사 로그"
        ],
        "JPA Auditing이 표준. Hibernate Envers (이력 관리)는 더 복잡한 요구사항에 적합"
    )

    create_class_detail_slide(prs,
        "QueryDslConfig",
        "com.hoji.config",
        "QueryDSL JPAQueryFactory 빈 등록",
        "타입 안전한 동적 쿼리 작성을 위한 QueryDSL 설정",
        [
            "컴파일 타임 쿼리 검증",
            "IDE 자동완성 지원",
            "복잡한 동적 쿼리 작성 용이",
            "JPQL보다 가독성 좋은 코드"
        ],
        "JOOQ, Criteria API도 있으나, QueryDSL이 Kotlin과 Spring 생태계에서 가장 인기 있음"
    )

    create_class_detail_slide(prs,
        "RedisConfig & RabbitMqConfig",
        "com.hoji.config",
        "Redis와 RabbitMQ 연결 및 설정",
        "캐싱, 세션 저장소(Redis) 및 비동기 메시징(RabbitMQ) 지원",
        [
            "RedisTemplate, StringRedisTemplate 빈 제공",
            "RabbitTemplate, MessageConverter 설정",
            "연결 풀 및 타임아웃 설정",
            "선택적 사용 (AutoConfiguration exclude)"
        ],
        "Redis: Lettuce(기본), Jedis. RabbitMQ: Kafka도 대안이나 용도가 다름"
    )

    create_class_detail_slide(prs,
        "WebClientConfig & RestTemplateConfig",
        "com.hoji.config",
        "HTTP 클라이언트 설정",
        "외부 API 호출을 위한 클라이언트 빈 제공",
        [
            "WebClient: 비동기, 논블로킹 (Reactive)",
            "RestTemplate: 동기, 블로킹 (전통적)",
            "타임아웃, 커넥션 풀 설정",
            "에러 핸들링 통합"
        ],
        "WebClient가 최신 표준. RestTemplate은 유지보수 모드이나 기존 코드 호환성에 유용"
    )

    create_class_detail_slide(prs,
        "AsyncConfig",
        "com.hoji.config",
        "@EnableAsync 및 ThreadPool 설정",
        "비동기 메서드 실행을 위한 ThreadPoolTaskExecutor 설정",
        [
            "@Async 어노테이션 지원",
            "ThreadPool 크기 최적화",
            "Graceful Shutdown 지원",
            "성능 향상 (I/O 대기 작업)"
        ],
        "Spring의 @Async가 간단. Kotlin Coroutines, Project Reactor도 대안이나 러닝 커브 존재"
    )

    create_class_detail_slide(prs,
        "CorsConfig",
        "com.hoji.config",
        "CORS(Cross-Origin Resource Sharing) 설정",
        "브라우저에서 다른 도메인의 API 호출 허용",
        [
            "허용 Origin, Method, Header 설정",
            "Credentials 지원",
            "Preflight 캐싱",
            "보안 정책 중앙 관리"
        ],
        "Spring의 CorsFilter 또는 @CrossOrigin 어노테이션. 전역 설정은 CorsFilter가 적합"
    )

    create_class_detail_slide(prs,
        "MetricsConfig",
        "com.hoji.config",
        "멀티 인스턴스 환경에서 메트릭 태깅",
        "각 인스턴스를 구분하기 위한 공통 태그 추가 (instance, host, port)",
        [
            "Prometheus가 인스턴스별로 메트릭 수집",
            "Grafana에서 인스턴스별/전체 조회",
            "Kubernetes Pod 이름 자동 주입",
            "운영 환경 모니터링 필수"
        ],
        "Micrometer의 표준 기능. 멀티 인스턴스 환경에서 필수적인 설정"
    )

    create_class_detail_slide(prs,
        "SwaggerConfig",
        "com.hoji.config",
        "SpringDoc OpenAPI 설정",
        "자동 API 문서 생성 및 Swagger UI 제공",
        [
            "코드에서 자동으로 문서 생성",
            "Swagger UI로 API 테스트",
            "JWT 인증 스키마 설정",
            "환경별 서버 정보"
        ],
        "SpringDoc (OpenAPI 3), Springfox (Swagger 2 - deprecated). SpringDoc이 최신 표준"
    )

    create_class_detail_slide(prs,
        "WebMvcConfig",
        "com.hoji.config",
        "Spring MVC 설정 (Interceptor 등록)",
        "RequestContext, Logging, Metrics Interceptor를 전역 등록",
        [
            "인터셉터 체인 구성",
            "요청 전/후 처리",
            "패턴별 인터셉터 적용",
            "CORS, MessageConverter 등 추가 설정 가능"
        ],
        "Spring MVC의 표준 설정 방식. WebMvcConfigurer 인터페이스 구현이 베스트 프랙티스"
    )

    create_class_detail_slide(prs,
        "LoggingProperties",
        "com.hoji.config.properties",
        "@ConfigurationProperties로 로깅 설정 외부화",
        "민감 정보 마스킹 대상 헤더를 application.yml에서 관리",
        [
            "설정 외부화로 유연성 향상",
            "환경별 다른 마스킹 정책",
            "타입 안전한 설정",
            "IDE 자동완성 지원"
        ],
        "@ConfigurationProperties가 Spring Boot 표준. @Value보다 타입 안전하고 구조화됨"
    )

    # 12. Domain
    create_section_header_slide(prs, "Domain: 도메인 엔티티")

    create_class_detail_slide(prs,
        "BaseEntity",
        "com.hoji.domain.common",
        "모든 엔티티가 상속받는 추상 클래스 (@MappedSuperclass)",
        "생성일시, 수정일시, 생성자, 수정자를 자동으로 관리",
        [
            "중복 코드 제거",
            "JPA Auditing과 통합",
            "@EntityListeners(AuditingEntityListener)",
            "모든 엔티티에 일관된 감사 필드"
        ],
        "JPA의 @MappedSuperclass가 표준. 대안으로 @Embeddable도 가능하나 상속이 더 직관적"
    )

    create_class_detail_slide(prs,
        "User (예제 엔티티)",
        "com.hoji.domain",
        "사용자 정보를 표현하는 JPA 엔티티",
        "베이스 프로젝트의 구조를 보여주는 예제 도메인",
        [
            "BaseEntity 상속으로 감사 필드 자동 관리",
            "@Entity, @Table 어노테이션",
            "컬럼 제약조건 명시",
            "Kotlin data class의 장점 활용"
        ],
        "실제 서비스에서는 비즈니스에 맞는 도메인으로 대체. User는 참고용 예제"
    )

    # 13. Repository, Service, Controller
    create_section_header_slide(prs, "Repository, Service, Controller")

    create_class_detail_slide(prs,
        "UserRepository",
        "com.hoji.repository",
        "Spring Data JPA Repository 인터페이스",
        "데이터 액세스 계층을 추상화하여 CRUD 자동 제공",
        [
            "기본 CRUD 메서드 자동 제공",
            "메서드 이름 기반 쿼리 생성",
            "QueryDSL과 조합 가능",
            "테스트 용이성"
        ],
        "JpaRepository가 표준. 복잡한 쿼리는 QueryDSL, @Query, 네이티브 쿼리 활용"
    )

    create_class_detail_slide(prs,
        "UserService",
        "com.hoji.service",
        "비즈니스 로직을 처리하는 서비스 계층",
        "컨트롤러와 레포지토리 사이에서 트랜잭션 및 비즈니스 규칙 처리",
        [
            "@Transactional로 트랜잭션 관리",
            "비즈니스 로직 캡슐화",
            "여러 레포지토리 조합",
            "테스트 가능한 구조"
        ],
        "전통적인 3-tier 아키텍처. DDD의 Domain Service, Application Service로 분리도 가능"
    )

    create_class_detail_slide(prs,
        "UserController",
        "com.hoji.controller",
        "REST API 엔드포인트를 제공하는 컨트롤러",
        "HTTP 요청을 받아 서비스 계층 호출 후 ApiResponse로 응답",
        [
            "@RestController로 JSON 자동 변환",
            "일관된 ApiResponse 응답",
            "DTO 변환으로 엔티티 노출 방지",
            "유효성 검증 (@Valid)"
        ],
        "Spring MVC 표준. GraphQL, gRPC 등도 대안이나 REST가 가장 보편적"
    )

    create_class_detail_slide(prs,
        "UserDto",
        "com.hoji.controller.dto",
        "API 요청/응답을 위한 데이터 전송 객체",
        "엔티티를 직접 노출하지 않고 필요한 정보만 전달",
        [
            "엔티티와 API 스펙 분리",
            "불필요한 정보 노출 방지",
            "유효성 검증 어노테이션",
            "버전별 API 스펙 관리 용이"
        ],
        "DTO 패턴이 표준. MapStruct로 변환 자동화 가능"
    )

    # 14. 운영 고려사항
    create_section_header_slide(prs, "운영 고려사항")

    create_content_slide(prs, "환경별 설정 관리", [
        {'text': 'Local, Dev, Prod 프로파일 분리', 'level': 0, 'font_size': 18},
        {'text': 'application.yml - 공통 설정', 'level': 1, 'font_size': 16},
        {'text': 'application-{profile}.yml - 환경별 차이', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Local: 디버깅 편의성', 'level': 0, 'font_size': 18},
        {'text': 'H2 콘솔, DEBUG 로그, 상세 에러 정보', 'level': 1, 'font_size': 14},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Dev: 개발 서버 특성', 'level': 0, 'font_size': 18},
        {'text': '적당한 로깅, 주요 Actuator 엔드포인트', 'level': 1, 'font_size': 14},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Prod: 보안 및 성능', 'level': 0, 'font_size': 18},
        {'text': '최소 로깅, 에러 정보 비노출, 압축, HTTP/2', 'level': 1, 'font_size': 14},
    ])

    create_content_slide(prs, "멀티 인스턴스 메트릭 관리", [
        {'text': 'Prometheus + Grafana 통합', 'level': 0, 'font_size': 18},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '각 인스턴스 식별', 'level': 0, 'font_size': 18},
        {'text': 'application, instance, host, port 태그', 'level': 1, 'font_size': 16},
        {'text': 'Kubernetes Pod 이름 자동 주입', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Kubernetes 배포', 'level': 0, 'font_size': 18},
        {'text': '3-replica 예시 (k8s/deployment.yml)', 'level': 1, 'font_size': 16},
        {'text': 'Prometheus 자동 스크래핑 annotation', 'level': 1, 'font_size': 16},
        {'text': 'Liveness/Readiness Probe', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Grafana 쿼리 예시', 'level': 0, 'font_size': 18},
        {'text': 'sum(api_calls_total) - 전체 합산', 'level': 1, 'font_size': 14},
        {'text': 'api_calls_total{instance="hoji-api-0"} - 특정 인스턴스', 'level': 1, 'font_size': 14},
    ])

    create_content_slide(prs, "Graceful Shutdown & 비동기", [
        {'text': 'Graceful Shutdown', 'level': 0, 'font_size': 18},
        {'text': 'server.shutdown: graceful', 'level': 1, 'font_size': 16},
        {'text': 'timeout-per-shutdown-phase: 30s', 'level': 1, 'font_size': 16},
        {'text': '진행 중인 요청 완료 후 종료', 'level': 1, 'font_size': 16},
        {'text': '무중단 배포 지원', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '비동기 처리 (@Async)', 'level': 0, 'font_size': 18},
        {'text': 'ThreadPool 최적화 (core: 5, max: 10)', 'level': 1, 'font_size': 16},
        {'text': '이메일, 알림 등 I/O 작업 비동기 처리', 'level': 1, 'font_size': 16},
        {'text': 'Graceful Shutdown과 통합', 'level': 1, 'font_size': 16},
        {'text': '응답 속도 향상', 'level': 1, 'font_size': 16},
    ])

    create_content_slide(prs, "보안 고려사항", [
        {'text': 'CORS 설정', 'level': 0, 'font_size': 18},
        {'text': '허용 Origin, Method 명시', 'level': 1, 'font_size': 16},
        {'text': 'Credentials 지원', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '민감 정보 마스킹', 'level': 0, 'font_size': 18},
        {'text': 'Authorization, Cookie 등 헤더 마스킹', 'level': 1, 'font_size': 16},
        {'text': '환경별 마스킹 정책 (LoggingProperties)', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'Prod 환경 보안 강화', 'level': 0, 'font_size': 18},
        {'text': '에러 스택트레이스 비노출', 'level': 1, 'font_size': 16},
        {'text': 'H2 콘솔 비활성화', 'level': 1, 'font_size': 16},
        {'text': '최소 Actuator 엔드포인트만 노출', 'level': 1, 'font_size': 16},
    ])

    # 15. 마무리
    create_section_header_slide(prs, "요약 및 결론")

    create_content_slide(prs, "베이스 프로젝트의 가치", [
        {'text': '일관성', 'level': 0, 'font_size': 20},
        {'text': '모든 서비스가 동일한 구조와 규칙 사용', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '생산성', 'level': 0, 'font_size': 20},
        {'text': '공통 기능 재구현 불필요', 'level': 1, 'font_size': 16},
        {'text': '새 서비스 시작 시간 단축', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '품질', 'level': 0, 'font_size': 20},
        {'text': '검증된 패턴과 베스트 프랙티스', 'level': 1, 'font_size': 16},
        {'text': '운영 환경을 고려한 설계', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '확장성', 'level': 0, 'font_size': 20},
        {'text': '멀티 인스턴스, 모니터링, 로깅 기반 마련', 'level': 1, 'font_size': 16},
    ])

    create_content_slide(prs, "향후 확장 가능성", [
        {'text': '인증/인가', 'level': 0, 'font_size': 18},
        {'text': 'Spring Security + JWT', 'level': 1, 'font_size': 16},
        {'text': 'OAuth 2.0, OIDC', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '분산 추적', 'level': 0, 'font_size': 18},
        {'text': 'Spring Cloud Sleuth + Zipkin', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '중앙 로그 관리', 'level': 0, 'font_size': 18},
        {'text': 'ELK Stack (Elasticsearch, Logstash, Kibana)', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': '데이터베이스 마이그레이션', 'level': 0, 'font_size': 18},
        {'text': 'Flyway 또는 Liquibase', 'level': 1, 'font_size': 16},
        {'text': '', 'level': 0, 'font_size': 10},
        {'text': 'API 버저닝', 'level': 0, 'font_size': 18},
        {'text': 'URL 버저닝, Header 버저닝', 'level': 1, 'font_size': 16},
    ])

    # 마지막 슬라이드
    create_title_slide(prs, "Q & A", "질문이 있으시면 언제든지 문의해 주세요")

    # 저장
    output_path = "/home/user/hoji/Hoji_Backend_Base_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"프리젠테이션이 생성되었습니다: {output_path}")
    print(f"총 {len(prs.slides)} 슬라이드")

if __name__ == "__main__":
    main()
